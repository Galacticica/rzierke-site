"""
File: build_slides.py
Author: Reagan Zierke <reaganzierke@gmail.com>
Date: 2026-02-03
Description: Uses python-pptx to build slides for ministry presentations.
"""

from __future__ import annotations

import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor

from ministry.models import Song  


BLANK_LINE_SPLIT = re.compile(r"\n\s*\n+")

FONT_NAME = "Yu Gothic UI Semilight"
TITLE_FONT_SIZE = 45
LYRICS_FONT_SIZE = 40

SLIDE_W = Inches(13.333) 
SLIDE_H = Inches(7.5)

MARGIN_X = Inches(0.75)
MARGIN_Y = Inches(0.6)


def _split_slides(text: str) -> list[str]:
    if not text:
        return []
    t = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    return [p.strip("\n") for p in BLANK_LINE_SPLIT.split(t) if p.strip()]


def _set_black_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)


def _add_centered_textbox(
    slide,
    text: str,
    *,
    font_name: str,
    font_size_pt: int,
) -> None:
    left = MARGIN_X
    top = MARGIN_Y
    width = SLIDE_W - (MARGIN_X * 2)
    height = SLIDE_H - (MARGIN_Y * 2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()

    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  

    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.rstrip()
        p.alignment = PP_ALIGN.CENTER

        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 1.0

        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = RGBColor(255, 255, 255)


def build_song_pptx_bytes(song: Song) -> bytes:
    prs = Presentation()

    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank_layout)
    _set_black_background(title_slide)
    _add_centered_textbox(
        title_slide,
        song.title,
        font_name=FONT_NAME,
        font_size_pt=TITLE_FONT_SIZE,
    )

    spacer_slide = prs.slides.add_slide(blank_layout)
    _set_black_background(spacer_slide)

    items = song.arrangement_items.select_related("section").all()

    for item in items:
        blocks = _split_slides(item.section.lyrics)

        for _ in range(item.repeat_count):
            for block in blocks:
                slide = prs.slides.add_slide(blank_layout)
                _set_black_background(slide)
                _add_centered_textbox(
                    slide,
                    block,
                    font_name=FONT_NAME,
                    font_size_pt=LYRICS_FONT_SIZE,
                )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
