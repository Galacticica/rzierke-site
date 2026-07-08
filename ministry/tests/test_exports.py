"""
File: test_exports.py
Description: Tests for the PPTX/PDF export helpers (build_slides, build_pdf)
and the export endpoints.
"""

import io

import pytest
from django.urls import reverse
from pptx import Presentation

from ministry.models import Song
from ministry.utils.build_pdf import build_song_print_pdf_bytes
from ministry.utils.build_slides import (
    _build_title_text,
    _split_slides,
    build_song_pptx_bytes,
)


class TestSplitSlides:
    def test_blank_line_separates_slides(self):
        assert _split_slides("line one\nline two\n\nline three") == [
            "line one\nline two",
            "line three",
        ]

    def test_crlf_normalized(self):
        assert _split_slides("first\r\n\r\nsecond") == ["first", "second"]

    def test_blank_lines_with_whitespace_still_split(self):
        assert _split_slides("first\n   \n\nsecond") == ["first", "second"]

    def test_empty_text_returns_no_slides(self):
        assert _split_slides("") == []
        assert _split_slides("   \n\n  ") == []


class TestBuildTitleText:
    def test_plain_title_without_lsb(self):
        song = Song(title="Amazing Grace")
        assert _build_title_text(song) == "Amazing Grace"

    def test_lsb_number_field_appended(self):
        song = Song(title="Amazing Grace", lsb_number="744")
        assert _build_title_text(song) == "Amazing Grace\nLSB 744"

    def test_lsb_extracted_from_title_when_field_missing(self):
        song = Song(title="Amazing Grace (LSB 744)")
        assert _build_title_text(song) == "Amazing Grace\nLSB 744"

    def test_lsb_field_takes_precedence_over_title(self):
        song = Song(title="Amazing Grace (LSB 123)", lsb_number="744")
        assert _build_title_text(song) == "Amazing Grace\nLSB 744"


def slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)


class TestBuildSongPptxBytes:
    def test_slide_count_and_title(self, song_with_arrangement):
        pptx_bytes = build_song_pptx_bytes(song_with_arrangement)
        prs = Presentation(io.BytesIO(pptx_bytes))
        slides = list(prs.slides)

        # title + spacer + (verse: 2 blocks x 1 repeat) + (chorus: 1 block x 2 repeats)
        assert len(slides) == 6

        title = slide_text(slides[0])
        assert "Amazing Grace" in title
        assert "LSB 744" in title

        # Spacer slide is empty.
        assert slide_text(slides[1]).strip() == ""

        lyric_texts = [slide_text(s) for s in slides[2:]]
        assert "Amazing grace how sweet the sound" in lyric_texts[0]
        assert "I once was lost but now am found" in lyric_texts[1]
        assert lyric_texts[2] == lyric_texts[3]  # chorus repeated
        assert "My chains are gone" in lyric_texts[2]


class TestBuildSongPrintPdfBytes:
    def test_returns_pdf_bytes(self, song_with_arrangement):
        pdf_bytes = build_song_print_pdf_bytes(song_with_arrangement)
        assert pdf_bytes.startswith(b"%PDF-")
        assert len(pdf_bytes) > 1000


class TestExportEndpoints:
    def test_pptx_export_response(self, client, song_with_arrangement):
        url = reverse("song-export-pptx", args=[song_with_arrangement.slug])
        response = client.get(url)

        assert response.status_code == 200
        assert (
            response["Content-Type"]
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        assert "amazing-grace.pptx" in response["Content-Disposition"]
        assert response.content[:2] == b"PK"  # zip container

    def test_pdf_export_response(self, client, song_with_arrangement):
        url = reverse("song-export-handout-pdf", args=[song_with_arrangement.slug])
        response = client.get(url)

        assert response.status_code == 200
        assert response["Content-Type"] == "application/pdf"
        assert "amazing-grace-handout.pdf" in response["Content-Disposition"]
        assert response.content.startswith(b"%PDF-")

    @pytest.mark.django_db
    def test_unknown_slug_404s_for_both(self, client):
        pptx_url = reverse("song-export-pptx", args=["no-such-song"])
        pdf_url = reverse("song-export-handout-pdf", args=["no-such-song"])

        assert client.get(pptx_url).status_code == 404
        assert client.get(pdf_url).status_code == 404
